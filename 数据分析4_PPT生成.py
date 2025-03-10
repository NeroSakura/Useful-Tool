import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from pathlib import Path
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 定义可能的工作表名称，根据实际数据调整
possible_sheet_names = {
    '智慧中国行-区域详情': ['智慧中国行-区域详情'],
    '客户研讨会-区域详情': ['客户研讨会-区域详情'],
    'AI科技品鉴会-区域详情': ['AI科技品鉴会-区域详情'],
    '创新之旅-区域详情': ['创新之旅-区域详情']
}

# 定义可能的列名，根据实际数据调整
possible_column_names = {
    '大区': ['大区'],
    '订单 $M': ['订单 $M']
}


def find_excel_files(folder_path):
    """
    在指定文件夹中查找 Excel 文件
    :param folder_path: 文件夹路径
    :return: 找到的 Excel 文件列表
    """
    excel_files = []
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.lower().endswith(('.xlsx', '.xls')):
                excel_files.append(os.path.join(root, file))
    return excel_files


def read_excel(file_path):
    """
    读取 Excel 文件中指定工作表的数据
    :param file_path: Excel 文件路径
    :return: 包含指定工作表数据的字典
    """
    data = {}
    xls = pd.ExcelFile(file_path)
    available_sheets = xls.sheet_names
    for target_name, alternatives in possible_sheet_names.items():
        for alt_name in alternatives:
            if alt_name in available_sheets:
                try:
                    df = xls.parse(alt_name)
                    data[target_name] = df
                    break
                except Exception as e:
                    print(f"读取工作表 '{alt_name}' 时出现错误: {e}")
        else:
            print(f"Excel 文件中未找到名为 '{target_name}' 的工作表，将跳过该工作表。")
    if not data:
        print("未找到任何有效的工作表，程序退出。")
        return None
    return data


def create_charts(slide, df, order_amount_col):
    """
    在幻灯片中创建饼图和柱状图
    :param slide: 幻灯片对象
    :param df: 数据 DataFrame
    :param order_amount_col: 订单金额列名
    """
    # 创建饼图
    pie_chart_data = ChartData()
    pie_chart_data.categories = df['大区']
    pie_chart_data.add_series('订单金额 $M', df[order_amount_col])

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(4)
    pie_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, pie_chart_data
    ).chart
    pie_chart.has_legend = True
    pie_chart.legend.include_in_layout = False
    pie_chart.series[0].data_labels.show_percentage = True
    pie_chart.series[0].data_labels.show_category_name = True

    # 创建柱状图
    bar_chart_data = CategoryChartData()
    bar_chart_data.categories = df['大区']
    bar_chart_data.add_series('订单金额 $M', df[order_amount_col])

    left = Inches(5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(4)
    bar_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, bar_chart_data
    ).chart
    bar_chart.has_legend = True
    bar_chart.legend.include_in_layout = False

    # 添加平均值线
    mean_value = df[order_amount_col].mean()
    value_axis = bar_chart.value_axis
    # 修改为设置 has_major_gridlines 属性来显示网格线
    value_axis.has_major_gridlines = True
    # 使用正确的枚举值设置虚线样式
    value_axis.major_gridlines.format.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # 手动计算平均值线的位置
    value_min = value_axis.minimum_scale if value_axis.minimum_scale is not None else 0
    value_max = value_axis.maximum_scale if value_axis.maximum_scale is not None else df[order_amount_col].max()
    y_mean = height * (1 - (mean_value - value_min) / (value_max - value_min))

    # 绘制平均值线
    line_start_x = left + Inches(0.1)
    line_start_y = top + y_mean
    line_end_x = left + width - Inches(0.1)
    line_end_y = top + y_mean
    line = slide.shapes.add_shape(
        MSO_SHAPE.LINE_INVERSE, line_start_x, line_start_y, line_end_x - line_start_x, line_end_y - line_start_y
    )
    line.line.color.rgb = RGBColor(255, 0, 0)
    line.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def create_ppt(data, output_path):
    """
    创建 PPT 文件
    :param data: 包含工作表数据的字典
    :param output_path: PPT 保存路径
    """
    prs = Presentation()
    title_layout = prs.slide_layouts[1]  # 使用标题和内容布局

    for sheet_name in data:
        df = data[sheet_name]
        # 查找实际的订单金额列名
        order_amount_col = None
        for alt_col in possible_column_names['订单 $M']:
            if alt_col in df.columns:
                order_amount_col = alt_col
                break
        if order_amount_col is None:
            print(f"工作表 '{sheet_name}' 中未找到订单金额列，将跳过该幻灯片。")
            continue

        slide = prs.slides.add_slide(title_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = sheet_name
            title.text_frame.paragraphs[0].font.size = Pt(18)
        else:
            print(f"无法为幻灯片 '{sheet_name}' 设置标题。")
            continue

        create_charts(slide, df, order_amount_col)

    prs.save(output_path)


def main():
    # 通过 input 获取文件路径
    path = input("请输入 Excel 文件的存放路径（可以是文件路径或文件夹路径）：")

    # 获取桌面路径
    desktop_path = Path.home() / "Desktop"
    output_ppt = desktop_path / "区域业务分析报告.pptx"

    if os.path.isfile(path) and path.lower().endswith(('.xlsx', '.xls')):
        excel_files = [path]
    elif os.path.isdir(path):
        excel_files = find_excel_files(path)
    else:
        print("输入的路径不是有效的 Excel 文件路径或文件夹路径，请检查后重新运行程序。")
        return

    all_data = {}
    for file_path in excel_files:
        data = read_excel(file_path)
        if data is None:
            continue
        all_data.update(data)

    create_ppt(all_data, output_ppt)
    print(f"PPT 已生成：{output_ppt}")


if __name__ == "__main__":
    main()